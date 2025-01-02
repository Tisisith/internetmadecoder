from pptx import Presentation

def create_bible_quiz_presentation():
    """
    Creates a PowerPoint presentation with Bible-related questions and answers.
    Each question slide is followed by an answer slide.
    """
    # Create a presentation object
    presentation = Presentation()

    # List of Bible-related questions, options, and answers
    questions = [
        {"question": "Who built the ark?", "options": ["A. Moses", "B. Noah", "C. Abraham", "D. David"], "answer": "B. Noah"},
        {"question": "Where was Jesus born?", "options": ["A. Jerusalem", "B. Bethlehem", "C. Nazareth", "D. Galilee"], "answer": "B. Bethlehem"},
        {"question": "What was Jesus' first miracle?", "options": ["A. Healing the blind", "B. Walking on water", "C. Turning water into wine", "D. Feeding the 5,000"], "answer": "C. Turning water into wine"},
        {"question": "Who killed Goliath?", "options": ["A. Saul", "B. David", "C. Solomon", "D. Samson"], "answer": "B. David"},
        {"question": "Who was swallowed by a great fish?", "options": ["A. Jonah", "B. Daniel", "C. Elijah", "D. Job"], "answer": "A. Jonah"},
        {"question": "What is the shortest verse in the Bible?", "options": ["A. 'God is love.'", "B. 'Jesus wept.'", "C. 'He is risen.'", "D. 'Rejoice always.'"], "answer": "B. 'Jesus wept.'"},
        {"question": "Who wrote most of the Psalms?", "options": ["A. Solomon", "B. David", "C. Moses", "D. Asaph"], "answer": "B. David"},
        {"question": "What did God create on the first day?", "options": ["A. Land", "B. Light", "C. Animals", "D. Plants"], "answer": "B. Light"},
        {"question": "Which book of the Bible contains the story of Samson and Delilah?","options": ["A. Judges", "B. Joshua", "C. Ruth", "D. 1 Samuel"], "answer": "A. Judges"},
        {"question": "What was the original language of most of the Old Testament?", "options": ["A. Greek", "B. Latin", "C. Hebrew", "D. Aramaic"], "answer": "C. Hebrew"},
        {"question": "Who was the first Christian martyr?","options": ["A. James", "B. Peter", "C. Stephen", "D. Paul"], "answer": "C. Stephen"},
        {"question": "What was the name of the man who replaced Judas Iscariot as an apostle?", "options": ["A. Barnabas", "B. Matthias", "C. Silas", "D. Timothy"], "answer": "B. Matthias"},
        {"question": "Which verse says 'The Lord bless you and keep you'?","options": ["A. Numbers 6:24", "B. Deuteronomy 6:4", "C. Joshua 1:9", "D. Exodus 20:12"],"answer": "A. Numbers 6:24"},
        {"question": "Which verse says 'The Lord is my light and my salvation'?","options": ["A. Psalm 27:1", "B. Psalm 46:1", "C. Psalm 91:1", "D. Psalm 121:1"], "answer": "A. Psalm 27:1"},
        {"question": "Where does the Bible first mention 'Satan' by name?", "options": ["A. Genesis 3:1", "B. Job 1:6", "C. Zechariah 3:1", "D. 1 Chronicles 21:1"], "answer": "D. 1 Chronicles 21:1"},
        {"question": "Which Psalm is the longest chapter in the Bible?","options": ["A. Psalm 119", "B. Psalm 117", "C. Psalm 23", "D. Psalm 91"],"answer": "A. Psalm 119"},
        {"question": "Where is the command 'Do not boil a young goat in its mother’s milk' found?","options": ["A. Exodus 23:19", "B. Leviticus 22:27", "C. Numbers 19:9", "D. Deuteronomy 14:21"], "answer": "A. Exodus 23:19"},
        {"question": "Where is it written, 'Can two walk together unless they agree'?","options": ["A. Amos 3:3", "B. Hosea 2:6", "C. Micah 6:8", "D. Joel 2:3"],"answer": "A. Amos 3:3"},

    ]

    # Add slides for each question and answer
    for item in questions:
        # Question slide
        question_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        question_slide.shapes.title.text = item["question"]
        content = "\n".join(item["options"])
        question_slide.placeholders[1].text = content

        # Answer slide
        answer_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        answer_slide.shapes.title.text = "Answer"
        answer_slide.placeholders[1].text = f"The correct answer is: {item['answer']}"

    # Save the presentation
    presentation.save("Bible_Quiz_Presentation.pptx")
    print("Presentation created and saved as 'Bible_Quiz_Presentation.pptx'")

# Run the function
create_bible_quiz_presentation()
